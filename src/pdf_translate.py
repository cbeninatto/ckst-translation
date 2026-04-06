import io
from collections import Counter
from typing import Callable, Dict, List, Optional, Tuple

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Pt

from .openai_translate import OpenAITranslator, TranslationItem, chunk_items
from .text_utils import apply_glossary_hard


EMU_PER_POINT = 12700


def _int_to_rgb(color_int: int) -> Tuple[int, int, int]:
    r = (color_int >> 16) & 255
    g = (color_int >> 8) & 255
    b = color_int & 255
    return r, g, b


def _pick_font_name(span_font_name: str) -> str:
    # Use common PowerPoint-safe fonts
    name = (span_font_name or "").lower()
    if "arial" in name:
        return "Arial"
    if "calibri" in name:
        return "Calibri"
    if "tahoma" in name:
        return "Tahoma"
    return "Arial"


def _has_letters(s: str) -> bool:
    return any(ch.isalpha() for ch in (s or ""))


def _extract_line_items(page: fitz.Page) -> List[dict]:
    """
    Extract line-level items from a PDF page:
    one item per visible text line with a union bbox.
    """
    d = page.get_text("dict")
    items: List[dict] = []

    for block in d.get("blocks", []):
        if block.get("type") != 0:
            continue

        for line in block.get("lines", []):
            spans = line.get("spans", [])
            if not spans:
                continue

            text = "".join(span.get("text", "") for span in spans).strip()
            if not text:
                continue
            if not _has_letters(text):
                continue

            rect = None
            max_size = 10.0
            font_name = "Arial"
            is_bold = False
            color = (0, 0, 0)

            for span in spans:
                bbox = span.get("bbox")
                if bbox:
                    srect = fitz.Rect(bbox)
                    rect = srect if rect is None else (rect | srect)

                size = float(span.get("size", 10.0))
                if size > max_size:
                    max_size = size

                span_font = span.get("font", "")
                font_name = _pick_font_name(span_font)
                if "bold" in span_font.lower():
                    is_bold = True

                color = _int_to_rgb(int(span.get("color", 0)))

            if rect is None:
                continue

            items.append(
                {
                    "text": text,
                    "rect": rect,
                    "font_size": max_size,
                    "font_name": font_name,
                    "bold": is_bold,
                    "color": color,
                }
            )

    return items


def _dominant_bg_color(img: Image.Image, box: Tuple[int, int, int, int], pad: int = 4) -> Tuple[int, int, int]:
    """
    Sample dominant color around a box from the rendered page image.
    This helps hide the original Portuguese text before placing English.
    """
    x0, y0, x1, y1 = box
    w, h = img.size

    x0 = max(0, x0 - pad)
    y0 = max(0, y0 - pad)
    x1 = min(w - 1, x1 + pad)
    y1 = min(h - 1, y1 + pad)

    px = img.load()
    samples: List[Tuple[int, int, int]] = []

    for x in range(x0, x1 + 1):
        samples.append(px[x, y0][:3] if isinstance(px[x, y0], tuple) else (255, 255, 255))
        samples.append(px[x, y1][:3] if isinstance(px[x, y1], tuple) else (255, 255, 255))

    for y in range(y0, y1 + 1):
        samples.append(px[x0, y][:3] if isinstance(px[x0, y], tuple) else (255, 255, 255))
        samples.append(px[x1, y][:3] if isinstance(px[x1, y], tuple) else (255, 255, 255))

    if not samples:
        return 255, 255, 255

    quantized = [((r // 16) * 16, (g // 16) * 16, (b // 16) * 16) for r, g, b in samples]
    return Counter(quantized).most_common(1)[0][0]


def _fit_font_size(text: str, box_w_pt: float, box_h_pt: float, start_size: float) -> float:
    """
    Basic best-effort font shrinker for PowerPoint text boxes.
    """
    size = max(8.0, min(start_size, 32.0))
    # rough width estimate: average char ≈ 0.5 * font_size points
    # rough line height: 1.2 * font_size
    for _ in range(18):
        est_chars_per_line = max(1, int(box_w_pt / max(1.0, size * 0.5)))
        est_lines = max(1, (len(text) // est_chars_per_line) + 1)
        est_height = est_lines * size * 1.2
        if est_height <= box_h_pt:
            return size
        size -= 0.75
        if size <= 7.0:
            return 7.0
    return max(7.0, size)


def translate_pdf_to_pptx_bytes(
    pdf_bytes: bytes,
    translator: OpenAITranslator,
    source_lang: str = "pt-BR",
    target_lang: str = "en",
    glossary: Optional[Dict[str, str]] = None,
    extra_instructions: str = "",
    on_progress: Optional[Callable[[str, int, int], None]] = None,
    render_scale: float = 2.0,
) -> bytes:
    """
    Workaround path:
    PDF -> translated PPTX

    Each PDF page becomes one PPTX slide with:
    - rendered page image as background
    - covered original text areas
    - translated text boxes placed over them
    """
    glossary = glossary or {}

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    prs = Presentation()

    if doc.page_count == 0:
        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()

    first_page = doc[0]
    base_w_pt = first_page.rect.width
    base_h_pt = first_page.rect.height

    prs.slide_width = int(base_w_pt * EMU_PER_POINT)
    prs.slide_height = int(base_h_pt * EMU_PER_POINT)

    total_pages = doc.page_count
    if on_progress:
        on_progress("pages", 0, total_pages)

    blank_layout = prs.slide_layouts[6]

    for page_idx in range(total_pages):
        page = doc[page_idx]
        if on_progress:
            on_progress("pages", page_idx + 1, total_pages)

        slide = prs.slides.add_slide(blank_layout)

        page_w_pt = page.rect.width
        page_h_pt = page.rect.height

        # Render page to image
        matrix = fitz.Matrix(render_scale, render_scale)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        page_img_bytes = pix.tobytes("png")
        pil_img = Image.open(io.BytesIO(page_img_bytes)).convert("RGBA")

        # Add background page image covering full slide
        slide.shapes.add_picture(
            io.BytesIO(page_img_bytes),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Extract PDF text lines
        line_items = _extract_line_items(page)
        if not line_items:
            continue

        # Translate text in chunks
        items = [TranslationItem(f"p{page_idx}_i{i}", li["text"]) for i, li in enumerate(line_items)]
        mapping: Dict[str, str] = {}

        total_items = len(items)
        done = 0
        if on_progress:
            on_progress("blocks", 0, max(1, total_items))

        for ch in chunk_items(items):
            mapping.update(
                translator.translate_batch(
                    ch,
                    source_lang=source_lang,
                    target_lang=target_lang,
                    glossary=glossary,
                    extra_instructions=extra_instructions,
                )
            )
            done += len(ch)
            if on_progress:
                on_progress("blocks", min(done, total_items), max(1, total_items))

        # Overlay translated text
        for i, li in enumerate(line_items):
            translated = mapping.get(f"p{page_idx}_i{i}", li["text"])
            translated = apply_glossary_hard(translated, glossary)

            rect: fitz.Rect = li["rect"]

            # Convert PDF points -> slide EMU
            left = int((rect.x0 / page_w_pt) * prs.slide_width)
            top = int((rect.y0 / page_h_pt) * prs.slide_height)
            width = int((rect.width / page_w_pt) * prs.slide_width)
            height = int((rect.height / page_h_pt) * prs.slide_height)

            if width <= 0 or height <= 0:
                continue

            # Sample dominant background color from rendered page image
            px_box = (
                int(rect.x0 * render_scale),
                int(rect.y0 * render_scale),
                int(rect.x1 * render_scale),
                int(rect.y1 * render_scale),
            )
            bg = _dominant_bg_color(pil_img, px_box, pad=5)

            # Add cover rectangle to hide Portuguese text
            cover = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                left,
                top,
                width,
                height,
            )
            cover.fill.solid()
            cover.fill.fore_color.rgb = __import__("pptx.dml.color").dml.color.RGBColor(bg[0], bg[1], bg[2])
            cover.line.fill.background()

            # Add translated text box
            tx = slide.shapes.add_textbox(left, top, width, height)
            tf = tx.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = translated

            font = run.font
            font.name = li["font_name"]
            font.bold = bool(li["bold"])

            box_w_pt = rect.width
            box_h_pt = rect.height
            fitted_size = _fit_font_size(translated, box_w_pt, box_h_pt, li["font_size"])
            font.size = Pt(fitted_size)

            r, g, b = li["color"]
            font.color.rgb = __import__("pptx.dml.color").dml.color.RGBColor(r, g, b)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
