import io
from typing import Callable, Dict, Optional

from pptx import Presentation

from .openai_translate import OpenAITranslator, TranslationItem, chunk_items
from .text_utils import apply_glossary_hard


def translate_pptx_bytes(
    pptx_bytes: bytes,
    translator: OpenAITranslator,
    source_lang: str = "pt-BR",
    target_lang: str = "en",
    glossary: Optional[Dict[str, str]] = None,
    extra_instructions: str = "",
    on_progress: Optional[Callable[[str, int, int], None]] = None,
) -> bytes:
    """
    Translates PPTX text in-place (same slides/shapes).
    """
    glossary = glossary or {}

    prs = Presentation(io.BytesIO(pptx_bytes))

    # Collect paragraph-level items
    items = []
    ptrs = []  # (item_id, paragraph, sample_run)

    total_slides = len(prs.slides)
    if on_progress:
        on_progress("slides", 0, max(1, total_slides))

    for si, slide in enumerate(prs.slides):
        if on_progress:
            on_progress("slides", si + 1, max(1, total_slides))

        for shape_index, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            if tf is None:
                continue

            for pi, p in enumerate(tf.paragraphs):
                text = "".join(run.text for run in p.runs) if p.runs else (p.text or "")
                if not text or not text.strip():
                    continue

                # skip pure placeholders/codes maybe? (translator protection already helps)
                item_id = f"s{si}_sh{shape_index}_p{pi}"
                items.append(TranslationItem(item_id, text))
                sample_run = p.runs[0] if p.runs else None
                ptrs.append((item_id, p, sample_run))

    # Translate in large chunks
    mapping: Dict[str, str] = {}
    total_items = len(items)
    done = 0
    if on_progress:
        on_progress("text", 0, max(1, total_items))

    for ch in chunk_items(items):
        mapping.update(
            translator.translate_batch(
                ch,
                source_lang=source_lang,
                target_lang=target_lang,
                glossary=glossary,
                extra_instructions=extra_instructions,
            )
        )
        done += len(ch)
        if on_progress:
            on_progress("text", min(done, total_items), max(1, total_items))

    # Write back
    for item_id, paragraph, sample_run in ptrs:
        new_text = mapping.get(item_id, None)
        if not new_text:
            continue
        new_text = apply_glossary_hard(new_text, glossary)

        # Preserve style: rewrite paragraph to a single run
        # This may simplify formatting, but keeps box positioning.
        # If you need strict per-run formatting, we can do a more advanced run mapping later.
        paragraph.text = new_text
        if sample_run is not None and paragraph.runs:
            r0 = paragraph.runs[0]
            try:
                r0.font.name = sample_run.font.name
                r0.font.size = sample_run.font.size
                r0.font.bold = sample_run.font.bold
                r0.font.italic = sample_run.font.italic
                r0.font.underline = sample_run.font.underline
                if sample_run.font.color and sample_run.font.color.rgb:
                    r0.font.color.rgb = sample_run.font.color.rgb
            except Exception:
                pass

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
